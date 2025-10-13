"""
CI構想プレゼンテーション自動生成テスト版
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import os
    print("✅ 必要なライブラリのインポートに成功しました")
    
    # 簡単なテスト用プレゼンテーション作成
    prs = Presentation()
    
    # タイトルスライド
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "CI構想テスト"
    subtitle.text = "プレゼンテーション自動生成テスト"
    
    # コンテンツスライド
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "テストスライド"
    content.text = "• テスト項目1\n• テスト項目2\n• テスト項目3"
    
    # 保存
    output_file = "CI_Test.pptx"
    prs.save(output_file)
    
    print(f"✅ テスト用プレゼンテーションを作成しました: {output_file}")
    print(f"📊 スライド数: {len(prs.slides)}枚")
    
    # ファイルが実際に作成されたか確認
    if os.path.exists(output_file):
        size = os.path.getsize(output_file)
        print(f"📁 ファイルサイズ: {size} bytes")
    else:
        print("❌ ファイルが作成されませんでした")

except ImportError as e:
    print(f"❌ インポートエラー: {e}")
    print("python-pptxライブラリがインストールされていない可能性があります")
except Exception as e:
    print(f"❌ エラーが発生しました: {e}")
    import traceback
    traceback.print_exc()